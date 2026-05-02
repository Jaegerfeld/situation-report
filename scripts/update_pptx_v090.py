# =============================================================================
# Autor:          Robert Seebauer
# Repository:     https://github.com/Jaegerfeld/situation-report
# KI-Unterstützung: Erstellt mit Unterstützung von Claude (Anthropic)
# Erstellt:       03.05.2026
# Geändert:       03.05.2026
# Lizenz:         BSD-3-Clause (siehe LICENSE)
#
# Fachliche Funktion:
#   Einmaliges Migrationsskript: aktualisiert die PPTX-Präsentationen auf v0.9.0.
#   Ändert Versionsnummern/Datum, aktualisiert ALPHA-Badge-Text und fügt neue
#   Folien für testdata_generator und helper in die Feature-Übersicht ein.
# =============================================================================

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

REPO = Path(__file__).parent.parent
FEATURES = REPO / "docs" / "situation_report_features.pptx"
PROZESS = REPO / "docs" / "situation_report_prozess.pptx"


def replace_in_slide(slide, old: str, new: str) -> None:
    """Replace a string in every run of every text frame on the slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def set_shape_text(slide, shape_name: str, new_text: str) -> bool:
    """
    Set the display text of a named shape, preserving first-run formatting.

    Clears all paragraphs beyond the first, then sets the first run's text.
    Returns True if the shape was found and updated.
    """
    for shape in slide.shapes:
        if shape.name != shape_name or not shape.has_text_frame:
            continue
        tf = shape.text_frame
        # Drop extra paragraphs
        while len(tf.paragraphs) > 1:
            tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ""
        else:
            from lxml import etree
            r = etree.SubElement(para._p, qn("a:r"))
            t = etree.SubElement(r, qn("a:t"))
            t.text = new_text
        return True
    return False


def copy_slide_append(prs: Presentation, src_idx: int):
    """
    Duplicate slide at src_idx and append it at the end of prs.

    Returns the new slide object.
    """
    src = prs.slides[src_idx]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    keep_tags = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    new_sp = new_slide.shapes._spTree
    src_sp = src.shapes._spTree

    for ch in list(new_sp):
        if ch.tag not in keep_tags:
            new_sp.remove(ch)
    for ch in src_sp:
        if ch.tag not in keep_tags:
            new_sp.append(copy.deepcopy(ch))

    return new_slide


# ---------------------------------------------------------------------------
# features.pptx
# ---------------------------------------------------------------------------

prs = Presentation(FEATURES)

# Slide 1: version + date
replace_in_slide(prs.slides[0], "v0.8.4", "v0.9.0")
replace_in_slide(prs.slides[0], "Stand April 2026", "Stand Mai 2026")

# Slide 2 (launcher): per-module badge description
set_shape_text(prs.slides[1], "TextBox 25", "BETA / ALPHA-Badges")
set_shape_text(
    prs.slides[1],
    "TextBox 26",
    "Orangefarbenes BETA-Badge im Titel zeigt den Gesamtreifegrad; "
    "jede Modul-Karte trägt ein eigenes ALPHA- oder BETA-Badge neben dem Namen",
)

# Slide 8 (Infrastruktur): version + test count
replace_in_slide(prs.slides[7], "aktuell v0.8.4", "aktuell v0.9.0")
replace_in_slide(prs.slides[7], "542 Tests", "628 Tests")

# --- New slide: testdata_generator (copy of slide 3 = index 2) ---
tdg = copy_slide_append(prs, 2)
set_shape_text(tdg, "TextBox 2", "testdata_generator")
for name, text in [
    ("TextBox 9",  "Synthetische Testdaten"),
    ("TextBox 10", "Erzeugt Jira-JSON-Exporte mit realistischem Workflow-Verlauf "
                   "— direkt mit transform_data verarbeitbar"),
    ("TextBox 13", "Konfigurierbar"),
    ("TextBox 14", "Projekt-Key, Issue-Anzahl, Zeitraum, Completion-Rate, "
                   "Backflow-Wahrscheinlichkeit und Seed — via GUI oder CLI"),
    ("TextBox 17", "Workflow-Datei"),
    ("TextBox 18", "Stages und Übergänge aus bestehender workflow.txt "
                   "— kompatibel mit transform_data-Format"),
    ("TextBox 21", "Reproduzierbar"),
    ("TextBox 22", "Seed-Parameter garantiert identische Ausgabe — "
                   "stabile, reproduzierbare Tests"),
    ("TextBox 25", "Issue-Typen"),
    ("TextBox 26", "Gewichtete Verteilung konfigurierbar "
                   "(z. B. 70 % Story, 20 % Bug, 10 % Task)"),
    ("TextBox 29", "GUI"),
    ("TextBox 30", "tkinter-Oberfläche mit allen Parametern; "
                   "Ladebalken nach 3 Sekunden"),
    ("TextBox 33", "CLI"),
    ("TextBox 34", "python -m testdata_generator; "
                   "alle Parameter als Kommandozeilenargumente"),
    ("TextBox 37", "Doppelklick-Starter"),
    ("TextBox 38", "testdata_generator_gui.pyw startet die GUI ohne Konsolenfenster"),
]:
    set_shape_text(tdg, name, text)

# --- New slide: helper (copy of slide 3 = index 2) ---
hlp = copy_slide_append(prs, 2)
set_shape_text(hlp, "TextBox 2", "helper — JSON Merger")
for name, text in [
    ("TextBox 9",  "Problem: Paginierung"),
    ("TextBox 10", "Jira REST API liefert max. 1000 Issues pro Abfrage "
                   "— bei größeren Projekten entstehen mehrere JSON-Dateien"),
    ("TextBox 13", "JSON Merger"),
    ("TextBox 14", "Fügt alle Teilexporte zu einem einzigen Jira-kompatiblen JSON zusammen "
                   "— direkt mit transform_data verarbeitbar"),
    ("TextBox 17", "Deduplizierung"),
    ("TextBox 18", "Issues mit identischer id werden standardmäßig nur einmal übernommen; "
                   "Warnung je Duplikat im Log"),
    ("TextBox 21", "Dedup-Kontrolle"),
    ("TextBox 22", "--no-dedup behält alle Issues "
                   "— nützlich wenn sich Zeitfenster der Abfragen überschneiden"),
    ("TextBox 25", "GUI"),
    ("TextBox 26", "Multi-File-Listbox + Ausgabepfad-Dialog + Dedup-Checkbox; "
                   "Ladebalken nach 3 Sekunden"),
    ("TextBox 29", "CLI"),
    ("TextBox 30", "python -m helper file1.json file2.json --output merged.json [--no-dedup]"),
    ("TextBox 33", "Ausgabeformat"),
    ("TextBox 34", "Gültiger Jira REST API Envelope (startAt, maxResults, total, issues) "
                   "— direkt mit transform_data verarbeitbar"),
    ("TextBox 37", "Doppelklick-Starter"),
    ("TextBox 38", "helper_gui.pyw startet die GUI ohne Konsolenfenster"),
]:
    set_shape_text(hlp, name, text)

prs.save(FEATURES)
print(f"features.pptx gespeichert ({len(prs.slides)} Folien)")

# ---------------------------------------------------------------------------
# prozess.pptx
# ---------------------------------------------------------------------------

prz = Presentation(PROZESS)

# Slide 1: date
replace_in_slide(prz.slides[0], "Stand April 2026", "Stand Mai 2026")

# Slide 3 (Daten bereitstellen): mention helper alongside transform_data
replace_in_slide(prz.slides[2], "Modul: transform_data", "Modul: transform_data, helper")

# Also update slide 2 (overview diagram) module annotation
replace_in_slide(prz.slides[1], "Modul: transform_data", "Modul: transform_data, helper")

prz.save(PROZESS)
print(f"prozess.pptx gespeichert ({len(prz.slides)} Folien)")
