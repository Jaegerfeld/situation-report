# =============================================================================
# Autor:          Robert Seebauer
# Repository:     https://github.com/Jaegerfeld/situation-report
# KI-Unterstützung: Erstellt mit Unterstützung von Claude (Anthropic)
# Erstellt:       03.05.2026
# Geändert:       23.05.2026
# Lizenz:         BSD-3-Clause (siehe LICENSE)
#
# Fachliche Funktion:
#   Erzeugt zwei zielgruppenspezifische PowerPoint-Präsentationen aus den
#   Inhalten der SituationReport-Dokumentation:
#   1. presentation_management.pptx – für Führungskräfte in der IT
#   2. presentation_rte_scrummaster.pptx – für RTEs und Scrum Master
# =============================================================================

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

REPO = Path(__file__).parent.parent
OUT_MGMT = REPO / "docs" / "presentation_management.pptx"
OUT_RTE  = REPO / "docs" / "presentation_rte_scrummaster.pptx"

IMG_GUI     = REPO / "docs" / "assets"
IMG_METRICS = REPO / "build_reports" / "prototype screens"

# ---------------------------------------------------------------------------
# Design constants
# ---------------------------------------------------------------------------

SW = 13.333   # slide width  [inches]
SH = 7.5      # slide height [inches]

C_DARK   = RGBColor(0x2C, 0x3E, 0x50)
C_BLUE   = RGBColor(0x29, 0x80, 0xB9)
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)
C_TEAL   = RGBColor(0x16, 0xA0, 0x85)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xEC, 0xF0, 0xF1)
C_SILVER = RGBColor(0xD5, 0xDB, 0xDB)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_TEXT   = RGBColor(0x2C, 0x3E, 0x50)

HEADER_H = 1.25   # standard header height
MARGIN   = 0.35   # left/right margin

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l: float, t: float, w: float, h: float,
          fill: RGBColor, line: RGBColor | None = None):
    """Add a solid-filled rectangle. Coordinates in inches."""
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def _txb(slide, l: float, t: float, w: float, h: float) -> object:
    """Add an empty text frame and return (textbox, text_frame)."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    return txb, tf


def _para(tf, text: str, size: float, color: RGBColor,
          bold: bool = False, italic: bool = False,
          align: PP_ALIGN = PP_ALIGN.LEFT,
          space_before: int = 0, first: bool = False) -> object:
    """Add a paragraph to an existing text frame. first=True uses paragraphs[0]."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return p


def _txt(slide, text: str, l: float, t: float, w: float, h: float,
         size: float, color: RGBColor,
         bold: bool = False, italic: bool = False,
         align: PP_ALIGN = PP_ALIGN.LEFT):
    """Shorthand: add a single-paragraph text box."""
    _, tf = _txb(slide, l, t, w, h)
    _para(tf, text, size, color, bold=bold, italic=italic, align=align, first=True)


def _image(slide, path: Path, l: float, t: float, w: float, h: float):
    """
    Insert a picture from `path` into the slot (l, t, w, h), preserving the
    image's aspect ratio and centering it within the slot. If the file does
    not exist, a gray placeholder rectangle with the filename is drawn so a
    missing asset is immediately visible without breaking the build.
    """
    if not path.exists():
        _rect(slide, l, t, w, h, C_SILVER, line=C_GRAY)
        _txt(slide, f"[missing: {path.name}]", l, t + h / 2 - 0.15, w, 0.3,
             11, C_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return
    with Image.open(path) as im:
        img_w, img_h = im.size
    img_ratio = img_w / img_h
    slot_ratio = w / h
    if img_ratio > slot_ratio:
        # image is wider than slot → fit width, center vertically
        new_w = w
        new_h = w / img_ratio
        new_l = l
        new_t = t + (h - new_h) / 2
    else:
        # image is taller than slot → fit height, center horizontally
        new_h = h
        new_w = h * img_ratio
        new_t = t
        new_l = l + (w - new_w) / 2
    slide.shapes.add_picture(str(path), Inches(new_l), Inches(new_t),
                             Inches(new_w), Inches(new_h))


def _bullets(slide, items: list, l: float, t: float, w: float, h: float,
             size: float = 14, color: RGBColor = C_TEXT,
             marker: str = "●", space: int = 4):
    """
    Add a list of bullet items.
    items: list of str  or  (str, level)  or  (str, level, color)
    """
    _, tf = _txb(slide, l, t, w, h)
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, level, clr = item, 0, color
        elif len(item) == 2:
            text, level = item; clr = color
        else:
            text, level, clr = item
        indent = "    " * level
        bullet = ("▸" if level > 0 else marker) + " "
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if space and i > 0:
            p.space_before = Pt(space)
        run = p.add_run()
        run.text = indent + bullet + text
        run.font.size = Pt(size)
        run.font.color.rgb = clr


# ---------------------------------------------------------------------------
# Mid-level helpers
# ---------------------------------------------------------------------------


def _header(slide, title: str, subtitle: str = "", color: RGBColor = C_DARK):
    """Standard dark header bar with white title + optional subtitle."""
    _rect(slide, 0, 0, SW, HEADER_H, color)
    _txt(slide, title, MARGIN, 0.1, SW - MARGIN * 2, 0.85, 26, C_WHITE, bold=True)
    if subtitle:
        _txt(slide, subtitle, MARGIN, 0.88, SW - MARGIN * 2, 0.32,
             12, C_LIGHT, italic=True)


def _badge(slide, text: str, l: float, t: float,
           bg: RGBColor = C_ORANGE, fg: RGBColor = C_WHITE, size: float = 13):
    """Small colored badge/pill (0.8" wide, 0.35" tall by default)."""
    _rect(slide, l, t, 0.9, 0.37, bg)
    _txt(slide, text, l + 0.05, t + 0.04, 0.8, 0.3, size, fg,
         bold=True, align=PP_ALIGN.CENTER)


def _card(slide, title: str, body: str, l: float, t: float, w: float, h: float,
          header_color: RGBColor = C_BLUE, body_bg: RGBColor = C_LIGHT):
    """A card with colored header and light body."""
    _rect(slide, l, t, w, 0.45, header_color)
    _txt(slide, title, l + 0.1, t + 0.06, w - 0.18, 0.35, 13,
         C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    _rect(slide, l, t + 0.45, w, h - 0.45, body_bg, line=C_SILVER)
    _txt(slide, body, l + 0.1, t + 0.53, w - 0.18, h - 0.62, 11,
         C_TEXT, align=PP_ALIGN.LEFT)


def _card_with_image(slide, title: str, img_path: Path, body: str,
                     l: float, t: float, w: float, h: float,
                     img_h: float,
                     header_color: RGBColor = C_BLUE,
                     body_bg: RGBColor = C_LIGHT):
    """Card with header, image at the top of the body, text below."""
    _rect(slide, l, t, w, 0.45, header_color)
    _txt(slide, title, l + 0.1, t + 0.06, w - 0.18, 0.35, 13,
         C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    _rect(slide, l, t + 0.45, w, h - 0.45, body_bg, line=C_SILVER)
    _image(slide, img_path, l + 0.1, t + 0.55, w - 0.2, img_h)
    txt_y = t + 0.55 + img_h + 0.08
    _txt(slide, body, l + 0.1, txt_y, w - 0.18, t + h - txt_y - 0.05,
         11, C_TEXT, align=PP_ALIGN.LEFT)


def _step_box(slide, number: str, label: str, desc: str,
              l: float, t: float, w: float, color: RGBColor = C_BLUE):
    """Numbered process step box."""
    h = 1.7
    _rect(slide, l, t, w, h, color)
    _txt(slide, number, l + 0.12, t + 0.1, 0.5, 0.5, 28, C_WHITE, bold=True)
    _txt(slide, label, l + 0.12, t + 0.62, w - 0.22, 0.4, 15, C_WHITE, bold=True)
    _txt(slide, desc, l + 0.12, t + 1.05, w - 0.22, 0.6, 11, C_LIGHT)


def _flow_step(slide, label: str, l: float, t: float,
               w: float = 2.1, h: float = 0.72, color: RGBColor = C_BLUE,
               optional: bool = False):
    """Compact flow step."""
    _rect(slide, l, t, w, h, color)
    _txt(slide, label, l + 0.08, t + 0.08, w - 0.14, h - 0.14, 14,
         C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    if optional:
        _txt(slide, "(optional)", l, t + h + 0.02, w, 0.25, 9,
             C_GRAY, italic=True, align=PP_ALIGN.CENTER)


def _arrow_right(slide, l: float, t: float):
    _txt(slide, "→", l, t, 0.4, 0.72, 22, C_DARK, bold=True,
         align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# ============================================================
# MANAGEMENT PRESENTATION
# ============================================================
# ---------------------------------------------------------------------------


def _m1_title(prs):
    """Slide 1: Title slide."""
    s = _blank(prs)
    _rect(s, 0, 0, SW, SH, C_DARK)                          # full background
    _rect(s, 0, 2.8, SW, 0.06, C_BLUE)                      # accent line
    _txt(s, "SituationReport", MARGIN, 1.2, SW - MARGIN * 2, 1.3,
         54, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, "Transparenz über den Entwicklungsfortschritt in Ihrer IT",
         MARGIN, 2.6, SW - MARGIN * 2, 0.6,
         22, C_LIGHT, align=PP_ALIGN.CENTER)
    _rect(s, 4.9, 3.15, 3.5, 0.5, C_ORANGE)
    _txt(s, "Für Führungskräfte in der IT",
         4.9, 3.2, 3.5, 0.45, 15, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, "Version 0.14.3  ·  Mai 2026",
         MARGIN, SH - 0.5, SW - MARGIN * 2, 0.4,
         11, C_GRAY, align=PP_ALIGN.CENTER)


def _m2_problem(prs):
    """Slide 2: Das Problem."""
    s = _blank(prs)
    _header(s, "Was wissen wir wirklich über unsere Lieferfähigkeit?")
    _bullets(s, [
        "Wie lange dauert es im Durchschnitt, bis eine Aufgabe erledigt ist?",
        "Wie viele Aufgaben liefert das Team pro Woche — und wird es mehr oder weniger?",
        "Wie viele Aufgaben sind gleichzeitig in Bearbeitung?",
        "Wo entstehen Engpässe im Entwicklungsprozess?",
        "Wie vorhersagbar sind unsere Liefertermine?",
    ], MARGIN, 1.45, 7.8, 4.2, size=16, color=C_TEXT, space=8)
    _rect(s, 8.6, 1.45, 4.4, 4.2, C_BLUE)
    _txt(s, "Ohne Messung\nkeine Antworten.",
         8.75, 1.7, 4.1, 1.2, 20, C_WHITE, bold=True)
    _txt(s,
         "Jira enthält alle Daten — "
         "aber sie werden selten systematisch ausgewertet. "
         "SituationReport macht diese Daten "
         "in Minuten sichtbar.",
         8.75, 3.1, 4.1, 2.3, 13, C_LIGHT)


def _m3_solution(prs):
    """Slide 3: Die Lösung."""
    s = _blank(prs)
    _header(s, "SituationReport — Kennzahlen direkt aus Jira")
    _txt(s, "Was es ist:", MARGIN, 1.45, 8.5, 0.35, 15, C_DARK, bold=True)
    _bullets(s, [
        "Eine lokale Werkzeugsammlung, die Jira-Exportdaten auswertet",
        "Erzeugt interaktive Diagramme und PDF-Berichte",
        "Keine Jira-Verbindung zur Laufzeit — einmaliger Export genügt",
    ], MARGIN, 1.82, 8.5, 1.5, size=14, color=C_TEXT, space=5)
    _txt(s, "Vier entscheidende Eigenschaften:", MARGIN, 3.45, 8.5, 0.35, 15, C_DARK, bold=True)
    props = [
        ("🔒  Lokal", "Läuft ausschließlich auf dem\neigenen Computer", C_BLUE),
        ("☁ Keine Cloud", "Keine Daten verlassen\ndas Unternehmen", C_TEAL),
        ("🆓  Kostenlos", "Open Source, BSD-Lizenz,\nkein Abo", C_GREEN),
        ("📊  Sofort nutzbar", "Download, entpacken,\nstarten", C_ORANGE),
    ]
    for i, (title, body, color) in enumerate(props):
        _card(s, title, body, MARGIN + i * 3.15, 3.85, 3.0, 1.5, color)


def _m4_workflow(prs):
    """Slide 4: Wie es funktioniert."""
    s = _blank(prs)
    _header(s, "Drei Schritte — vom Export zum Bericht")
    steps = [
        ("1", "Jira exportieren", "Jira-Board als JSON-Datei exportieren. Bei großen Projekten ggf. mehrere Dateien — der Helper fügt sie zusammen.", C_BLUE),
        ("2", "Transform Data", "Die Exportdatei einlesen. Das Tool berechnet, wie viel Zeit jede Aufgabe in welcher Bearbeitungsphase verbracht hat.", C_TEAL),
        ("3", "Build Reports", "Mit einem Klick werden alle Kennzahlen als interaktive Diagramme im Browser angezeigt oder als PDF exportiert.", C_GREEN),
    ]
    for i, (num, label, desc, color) in enumerate(steps):
        x = MARGIN + i * 4.3
        _step_box(s, num, label, desc, x, 1.5, 4.0, color)
        if i < 2:
            _txt(s, "→", x + 4.1, 2.1, 0.35, 0.6, 26, C_DARK, bold=True, align=PP_ALIGN.CENTER)
    _rect(s, MARGIN, 3.35, SW - MARGIN * 2, 0.06, C_SILVER)
    _txt(s, "Keine Installation von Abhängigkeiten nötig — das portable Paket enthält alles.",
         MARGIN, 3.5, SW - MARGIN * 2, 0.5, 13, C_GRAY,
         italic=True, align=PP_ALIGN.CENTER)


def _m5_metrics(prs):
    """Slide 5: Die vier Kernmetriken."""
    s = _blank(prs)
    _header(s, "Vier Kernfragen — vier Metriken")
    metrics = [
        ("Flow Time  /  Cycle Time",
         "Wie lange dauert eine Aufgabe?\n\nVon Bearbeitungsbeginn bis Abschluss. "
         "Zeigt Vorhersagbarkeit und Engpässe.", C_BLUE),
        ("Flow Velocity  /  Throughput",
         "Wie viele Aufgaben werden fertig?\n\nPro Tag, Woche und Planning Interval. "
         "Zeigt Kapazität und Trends.", C_TEAL),
        ("Flow Load  /  WIP",
         "Was ist gerade in Bearbeitung?\n\nAktuelle Aufgaben pro Phase, mit Alter. "
         "Zeigt Überlast und Blockierungen.", C_ORANGE),
        ("Cumulative Flow Diagram",
         "Wie fließt die Arbeit durchs System?\n\nZeigt Engpässe, Warteschlagen und "
         "ob In-flow und Out-flow im Gleichgewicht sind.", C_GREEN),
    ]
    for i, (title, body, color) in enumerate(metrics):
        row, col = divmod(i, 2)
        x = MARGIN + col * 6.4
        y = 1.45 + row * 2.85
        _card(s, title, body, x, y, 6.2, 2.6, color)


def _m6_status(prs):
    """Slide 6: Aktueller Stand."""
    s = _blank(prs)
    _header(s, "Heute verfügbar — Version 0.14.3")
    _txt(s, "Verfügbare Module:", MARGIN, 1.4, 8.0, 0.35, 15, C_DARK, bold=True)
    modules = [
        ("Transform Data", "BETA", C_ORANGE, "Jira-Rohdaten aufbereiten — Kernmodul"),
        ("Build Reports",  "BETA", C_ORANGE, "Flow-Metriken und PDF-Berichte — Kernmodul"),
        ("Testdata Generator", "BETA", C_ORANGE, "Synthetische Testdaten erzeugen"),
        ("Helper",         "ALPHA", C_RED,   "Mehrere JSON-Exporte zusammenführen"),
    ]
    for i, (name, badge, badge_color, desc) in enumerate(modules):
        y = 1.85 + i * 0.68
        _rect(s, MARGIN, y, 8.0, 0.55, C_LIGHT, line=C_SILVER)
        _badge(s, badge, MARGIN + 0.1, y + 0.09, bg=badge_color, size=10)
        _txt(s, name, MARGIN + 1.12, y + 0.08, 2.5, 0.38, 14, C_DARK, bold=True)
        _txt(s, desc,  MARGIN + 3.7,  y + 0.1,  4.5, 0.35, 13, C_GRAY)
    _rect(s, MARGIN, 4.7, SW - MARGIN * 2, 0.06, C_SILVER)
    plat_items = [
        "Windows (ZIP, enthält Python & Chrome — keine Installation nötig)",
        "macOS Apple Silicon (ZIP, Python-Umgebung wird beim ersten Start eingerichtet)",
        "Linux x64 (ZIP, Python-Umgebung wird beim ersten Start eingerichtet)",
    ]
    _txt(s, "Plattformen:", MARGIN, 4.85, 5.0, 0.35, 15, C_DARK, bold=True)
    _bullets(s, plat_items, MARGIN, 5.25, SW - MARGIN * 2, 1.8, size=12, space=3)


def _m7_roadmap(prs):
    """Slide 7: Roadmap."""
    s = _blank(prs)
    _header(s, "Roadmap — Was noch kommt")
    _rect(s, MARGIN, 1.45, SW - MARGIN * 2, 0.06, C_SILVER)
    now_items = [
        ("Transform Data  [BETA]", 0, C_TEXT),
        ("Build Reports  [BETA]", 0, C_TEXT),
        ("Testdata Generator  [BETA]", 0, C_TEXT),
        ("Helper  [ALPHA]", 0, C_TEXT),
    ]
    next_items = [
        ("Get Data  — Direktabruf aus Jira\n(kein manueller Export mehr nötig)", 0, C_TEXT),
    ]
    later_items = [
        ("Simulate  — Prognosen und Forecasting\n(Wann sind wir fertig? Wahrscheinlichkeits-basiert)", 0, C_TEXT),
    ]
    cols = [
        ("Jetzt verfügbar", now_items, C_GREEN),
        ("Als nächstes", next_items, C_BLUE),
        ("Mittelfristig", later_items, C_ORANGE),
    ]
    for i, (label, items, color) in enumerate(cols):
        x = MARGIN + i * 4.3
        _rect(s, x, 1.6, 4.1, 0.52, color)
        _txt(s, label, x + 0.12, 1.65, 3.9, 0.45, 16, C_WHITE, bold=True)
        _bullets(s, items, x + 0.12, 2.2, 3.9, 4.5, size=13, color=C_TEXT, space=6)


def _m8_next(prs):
    """Slide 8: Nächste Schritte."""
    s = _blank(prs)
    _header(s, "Jetzt starten")
    steps = [
        ("⬇  Herunterladen",
         "github.com/Jaegerfeld/situation-report/releases\n"
         "→ Neueste stabile Version als ZIP herunterladen", C_BLUE),
        ("🧪  Ausprobieren",
         "Testdata Generator starten → synthetische Daten erzeugen\n"
         "→ Transform Data → Build Reports\n"
         "Keine echten Daten nötig für den ersten Test.", C_TEAL),
        ("📂  Mit eigenen Daten",
         "Jira-Board als JSON exportieren\n"
         "→ Workflow-Datei erstellen (Stationen definieren)\n"
         "→ Transform Data → Build Reports", C_GREEN),
        ("💬  Feedback",
         "github.com/Jaegerfeld/situation-report/issues\n"
         "Fehler melden, Wünsche äußern, Beitrag leisten", C_ORANGE),
    ]
    for i, (title, body, color) in enumerate(steps):
        row, col = divmod(i, 2)
        x = MARGIN + col * 6.4
        y = 1.45 + row * 2.75
        _card(s, title, body, x, y, 6.2, 2.5, color)


def build_management(path: Path) -> None:
    """Build the management presentation and save to path."""
    prs = _new_prs()
    _m1_title(prs)
    _m2_problem(prs)
    _m3_solution(prs)
    _m4_workflow(prs)
    _m5_metrics(prs)
    _m6_status(prs)
    _m7_roadmap(prs)
    _m8_next(prs)
    prs.save(path)
    print(f"Management PPTX gespeichert ({len(prs.slides)} Folien): {path.name}")


# ---------------------------------------------------------------------------
# ============================================================
# RTE / SCRUM MASTER PRESENTATION
# ============================================================
# ---------------------------------------------------------------------------


def _r1_title(prs):
    s = _blank(prs)
    _rect(s, 0, 0, SW, SH, C_DARK)
    _rect(s, 0, 2.8, SW, 0.06, C_TEAL)
    _txt(s, "SituationReport", MARGIN, 1.2, SW - MARGIN * 2, 1.3,
         54, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, "Flow-Metriken für Release Train Engineers und Scrum Master",
         MARGIN, 2.6, SW - MARGIN * 2, 0.65,
         20, C_LIGHT, align=PP_ALIGN.CENTER)
    _rect(s, 3.9, 3.15, 5.5, 0.5, C_TEAL)
    _txt(s, "Praktische Anleitung  ·  Alle Metriken erklärt",
         3.9, 3.2, 5.5, 0.45, 14, C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, "Version 0.14.3  ·  Mai 2026",
         MARGIN, SH - 0.5, SW - MARGIN * 2, 0.4,
         11, C_GRAY, align=PP_ALIGN.CENTER)


def _r_launcher(prs):
    """Slide: Launcher screenshot + module status + UI language."""
    s = _blank(prs)
    _header(s, "Der Launcher — Einstiegspunkt für alle Module", color=C_TEAL)
    _image(s, IMG_GUI / "Launcher-GUI.png", MARGIN, 1.5, 6.6, 5.4)
    _txt(s, "Vier Module heute verfügbar:",
         7.3, 1.5, 5.7, 0.35, 14, C_DARK, bold=True)
    mods = [
        ("Transform Data",     "BETA",  C_ORANGE),
        ("Build Reports",      "BETA",  C_ORANGE),
        ("Testdata Generator", "BETA",  C_ORANGE),
        ("Helper",             "ALPHA", C_RED),
    ]
    for i, (name, badge, badge_color) in enumerate(mods):
        y = 1.95 + i * 0.55
        _rect(s, 7.3, y, 5.7, 0.48, C_LIGHT, line=C_SILVER)
        _badge(s, badge, 7.4, y + 0.06, bg=badge_color, size=10)
        _txt(s, name, 8.42, y + 0.08, 4.3, 0.35, 13, C_DARK, bold=True)
    _rect(s, 7.3, 4.4, 5.7, 0.06, C_SILVER)
    _txt(s, "Oberflächensprache (Flagge oben rechts):",
         7.3, 4.55, 5.7, 0.35, 13, C_DARK, bold=True)
    _bullets(s, [
        "Klick schaltet durch 5 Sprachen: DE · EN · RO · PT · FR",
        "Auswahl wird in prefs.json gespeichert",
        "Englisch ist seit v0.13.0 die Standard-Sprache",
    ], 7.3, 4.9, 5.7, 2.0, size=12, color=C_TEXT, space=4)


def _r2_workflow(prs):
    s = _blank(prs)
    _header(s, "Der Ablauf — vom Jira-Export zum Bericht", color=C_TEAL)
    # Flow diagram
    steps = [
        ("📥 Jira\nExport", C_BLUE, False),
        ("🔧 Helper", C_ORANGE, True),
        ("🔄 Transform\nData", C_TEAL, False),
        ("📊 Build\nReports", C_GREEN, False),
    ]
    total_w = len(steps) * 2.2 + (len(steps) - 1) * 0.5
    start_x = (SW - total_w) / 2
    y_flow = 1.6
    for i, (label, color, optional) in enumerate(steps):
        x = start_x + i * 2.7
        _flow_step(s, label, x, y_flow, w=2.2, h=0.9, color=color, optional=optional)
        if i < len(steps) - 1:
            _arrow_right(s, x + 2.25, y_flow + 0.09)
    # Descriptions
    descs = [
        ("Jira → Board als JSON exportieren.\nBei > 1.000 Issues: mehrere\nDateien nötig.", start_x),
        ("Nur bei mehreren\nExportdateien:\nzusammenführen.", start_x + 2.7),
        ("JSON + Workflow-Datei\neinlesen → drei\nExcel-Tabellen erzeugen.", start_x + 5.4),
        ("Excel-Tabellen laden\n→ Diagramme und PDF\nerzeugen.", start_x + 8.1),
    ]
    for text, x in descs:
        _txt(s, text, x, y_flow + 1.15, 2.2, 1.5, 12, C_GRAY, italic=True)
    _rect(s, MARGIN, 3.4, SW - MARGIN * 2, 0.06, C_SILVER)
    _txt(s,
         "Tipp: Testdata Generator nutzt 'Create Report' (v0.12.0) und "
         "erzeugt Schritt 3+4 in einem Klick — ideal fuer den ersten Test.",
         MARGIN, 3.55, SW - MARGIN * 2, 0.55, 13, C_BLUE, italic=True)
    _txt(s, "Beispiel-Workflow-Datei (eine Stage pro Zeile):",
         MARGIN, 4.2, 6.0, 0.35, 13, C_DARK, bold=True)
    wf_lines = [
        "Funnel:New:Open",
        "Analysis:In Analysis",
        "Implementation:In Progress:In Review",
        "Done:Closed",
        "<First>Analysis",
        "<InProgress>Implementation",
        "<Closed>Done",
    ]
    _rect(s, MARGIN, 4.58, 5.5, 2.6, C_LIGHT, line=C_SILVER)
    _, tf = _txb(s, MARGIN + 0.15, 4.68, 5.2, 2.4)
    for i, line in enumerate(wf_lines):
        _para(tf, line, 12, C_TEXT, italic=line.startswith("<"), first=(i == 0), space_before=3 if i > 0 else 0)
    _txt(s,
         "Stage:Alias — mehrere Jira-Status auf eine Stage mappen.\n"
         "<First>/<InProgress>/<Closed> — markieren Beginn, aktive Phase\n"
         "und Ende der Messung.\n\n"
         "(Quelle: tests/testdata/fixtures/workflow_simple.txt)",
         MARGIN + 5.7, 4.58, 7.0, 2.6, 12, C_GRAY)


def _r3_transform(prs):
    s = _blank(prs)
    _header(s, "Transform Data — Rohdaten in strukturierte Tabellen", color=C_TEAL)
    # Input/Output columns
    _txt(s, "Eingaben:", MARGIN, 1.45, 3.8, 0.35, 15, C_DARK, bold=True)
    _card(s, "📄  Jira-JSON-Export",
          "Die aus Jira exportierte Datei mit allen Issues und ihrer Änderungshistorie.",
          MARGIN, 1.85, 3.8, 1.1, C_BLUE, C_LIGHT)
    _card(s, "📝  Workflow-Datei",
          "Textdatei mit den Stage-Namen des Projekts und der Markierung für Start/Ende der Messung.",
          MARGIN, 3.05, 3.8, 1.1, C_BLUE, C_LIGHT)
    _txt(s, "Ausgaben:", 4.5, 1.45, 8.5, 0.35, 15, C_DARK, bold=True)
    outputs = [
        ("IssueTimes.xlsx",
         "Eine Zeile pro Aufgabe.\nEnthält: Erstellt / First Date / Closed Date sowie "
         "die Zeit (in Minuten) in jeder Stage.\nGrundlage für alle Build-Reports-Metriken.", C_GREEN),
        ("Transitions.xlsx",
         "Vollständiges Statuswechsel-Protokoll.\nJede Statusänderung jeder Aufgabe "
         "mit Zeitstempel.\nNützlich für manuelle Analysen.", C_TEAL),
        ("CFD.xlsx",
         "Tägliche Eintrittszählungen je Stage.\nGrundlage für das\nCumulative Flow Diagram.", C_ORANGE),
    ]
    for i, (title, body, color) in enumerate(outputs):
        _card(s, title, body, 4.5, 1.85 + i * 1.8, 8.5, 1.6, color)
    # Hand-over hint (under the Eingaben column)
    _rect(s, MARGIN, 4.35, 3.8, 2.7, C_LIGHT, line=C_SILVER)
    _txt(s, "🤝  Hand-over (v0.14.0)", MARGIN + 0.15, 4.45, 3.5, 0.35,
         13, C_DARK, bold=True)
    _txt(s,
         "Nach erfolgreichem Lauf öffnet der\n"
         "'Open in build_reports'-Button das\n"
         "nächste Modul mit allen drei XLSX +\n"
         "Workflow bereits vorbefüllt.\n\n"
         "Details: siehe nächste Folie.",
         MARGIN + 0.15, 4.8, 3.5, 2.2, 11, C_GRAY)


def _r_handover(prs):
    """Slide: Hand-over transform_data -> build_reports (v0.14.0)."""
    s = _blank(prs)
    _header(s, "Hand-over: transform_data → build_reports (v0.14.0)",
            color=C_TEAL)
    _txt(s, "Einbahnstraßen-Übergabe in drei Schritten:",
         MARGIN, 1.42, SW - MARGIN * 2, 0.35, 14, C_GRAY, italic=True)
    steps = [
        ("1", "Transform läuft",
         "transform_data erzeugt IssueTimes,\nTransitions und CFD wie gewohnt.",
         C_BLUE),
        ("2", "Open in build_reports",
         "Neuer Button schreibt die drei XLSX\n"
         "+ Workflow in eine temporäre Template-\n"
         "Datei und startet build_reports.",
         C_TEAL),
        ("3", "Vorbefüllt starten",
         "build_reports öffnet sich mit allen\n"
         "Dateien, PI-Konfig, Filtern und\n"
         "Metrik-Auswahl bereits gesetzt.",
         C_GREEN),
    ]
    for i, (num, label, desc, color) in enumerate(steps):
        x = MARGIN + i * 4.3
        _step_box(s, num, label, desc, x, 1.95, 4.0, color)
        if i < 2:
            _txt(s, "→", x + 4.1, 2.55, 0.35, 0.6, 26, C_DARK, bold=True,
                 align=PP_ALIGN.CENTER)
    _txt(s, "CLI-Äquivalent für eigene Skripte:",
         MARGIN, 4.05, SW - MARGIN * 2, 0.35, 13, C_DARK, bold=True)
    _rect(s, MARGIN, 4.42, SW - MARGIN * 2, 0.7, C_LIGHT, line=C_SILVER)
    _txt(s, "python -m build_reports --gui-template handover.json",
         MARGIN + 0.2, 4.55, SW - MARGIN * 2 - 0.4, 0.5,
         15, C_DARK, bold=True)
    _rect(s, MARGIN, 5.4, SW - MARGIN * 2, 0.06, C_SILVER)
    _bullets(s, [
        "Aktuell einseitige Hand-over (transform_data → build_reports).",
        "Wenn in transform_data ein Project-Template geladen ist, werden "
        "dessen build_reports-Einstellungen (PI, Filter, Metriken) mit "
        "übergeben.",
        "Weitere Modul-Paare folgen — siehe Roadmap.",
    ], MARGIN, 5.55, SW - MARGIN * 2, 1.7, size=12, color=C_TEXT,
         marker="💡", space=5)


def _r4_build_intro(prs):
    s = _blank(prs)
    _header(s, "Build Reports — Fünf Metriken, ein Klick", color=C_TEAL)
    _txt(s, "IssueTimes.xlsx laden → Metriken auswählen → Berichte erstellen",
         MARGIN, 1.42, SW - MARGIN * 2, 0.4, 14, C_GRAY, italic=True)
    metrics = [
        ("Flow Time  /  Cycle Time", "Durchlaufzeit je Aufgabe (Boxplot + Scatterplot mit Trend)", C_BLUE),
        ("Flow Velocity  /  Throughput", "Abschlüsse pro Tag, Woche und Planning Interval", C_TEAL),
        ("Flow Load  /  WIP", "Aktuelle offene Aufgaben je Stage mit Alter", C_ORANGE),
        ("Cumulative Flow Diagram", "Kumulierter Fluss aller Issues über Zeit", C_GREEN),
        ("Flow Distribution", "Verteilung nach Typ, Stage-Dominanz, CT je Typ", C_RED),
    ]
    for i, (title, desc, color) in enumerate(metrics):
        y = 1.9 + i * 0.96
        _rect(s, MARGIN, y, SW - MARGIN * 2, 0.82, C_LIGHT, line=C_SILVER)
        _rect(s, MARGIN, y, 0.18, 0.82, color)
        _txt(s, title, MARGIN + 0.3, y + 0.06, 5.5, 0.35, 14, C_DARK, bold=True)
        _txt(s, desc,  MARGIN + 0.3, y + 0.42, 12.0, 0.35, 12, C_GRAY)
    _rect(s, MARGIN, 6.8, SW - MARGIN * 2, 0.4, C_LIGHT, line=C_SILVER)
    _txt(s,
         "Terminologie umschaltbar: SAFe (Flow Time, Flow Velocity, Flow Load) "
         "↔ Global (Cycle Time, Throughput, WIP)",
         MARGIN + 0.2, 6.85, SW - MARGIN * 2 - 0.3, 0.32, 12, C_GRAY)


def _r5_flow_time(prs):
    s = _blank(prs)
    _header(s, "Flow Time / Cycle Time — Wie lange dauert eine Aufgabe?", color=C_BLUE)
    _txt(s, "Messung:", MARGIN, 1.45, 3.0, 0.3, 13, C_DARK, bold=True)
    _txt(s, "Von First Date (erste Aktivität) bis Closed Date (Abschluss) — in Kalendertagen.",
         MARGIN, 1.78, SW - MARGIN * 2, 0.35, 13, C_TEXT)
    # Two chart cards with embedded screenshots
    _card_with_image(s, "📦  Boxplot — Verteilung",
          IMG_METRICS / "Flow_time_boxplot.png",
          "Enge Box = vorhersagbar. Breite Box = unvorhersagbar.\n"
          "Header-Statistiken: Min / Q1 / Median / Ø / Q3 / Max.\n"
          "90d CT% = Anteil der Aufgaben ≤ 90 Tage.",
          MARGIN, 2.25, 6.2, 3.0, img_h=1.55, header_color=C_BLUE)
    _card_with_image(s, "⚫  Scatterplot — Verlauf über Zeit",
          IMG_METRICS / "Flow_time_scatterplot.png",
          "Jeder Punkt = eine abgeschlossene Aufgabe.\n"
          "Referenzlinien: Median, P85, P95. Trendlinie zeigt, "
          "ob das Team schneller oder langsamer wird.",
          6.8, 2.25, 6.2, 3.0, img_h=1.55, header_color=C_BLUE)
    _rect(s, MARGIN, 5.4, SW - MARGIN * 2, 0.06, C_SILVER)
    _bullets(s, [
        "P85-Linie (85. Perzentil) = 85 % aller Aufgaben waren schneller als dieser Wert.",
        "Faustregel: P85 als Service Level Expectation (SLE) kommunizieren.",
        "Trendlinie zeigt Verbesserung/Verschlechterung über Zeit — früher erkennbar als Durchschnittswerte.",
    ], MARGIN, 5.55, SW - MARGIN * 2, 1.7, size=12, color=C_TEXT, marker="💡", space=4)


def _r6_flow_velocity(prs):
    s = _blank(prs)
    _header(s, "Flow Velocity / Throughput — Wie viele Aufgaben werden fertig?", color=C_TEAL)
    _txt(s, "Messung: Anzahl abgeschlossener Aufgaben (mit Closed Date) pro Zeitraum.",
         MARGIN, 1.42, SW - MARGIN * 2, 0.35, 13, C_TEXT)
    charts = [
        ("📊  Tagesfrequenz",
         IMG_METRICS / "Flow_velocity_daily.png",
         "Histogramm: wie viele Aufgaben werden typischerweise an einem Tag abgeschlossen?",
         C_TEAL),
        ("📈  Wochenverlauf",
         IMG_METRICS / "flow_velocity_weekly.png",
         "Linienchart: zeigt Schwankungen und Trends über Wochen und Quartale.",
         C_BLUE),
        ("📅  PI-Verlauf",
         IMG_METRICS / "Flow_velocity_ProgramIncrement.png",
         "Balkendiagramm pro PI mit Durchschnittslinie — ideal für PI-Retro und PI-Planning.",
         C_ORANGE),
    ]
    for i, (title, img, body, color) in enumerate(charts):
        x = MARGIN + i * 4.25
        _card_with_image(s, title, img, body, x, 1.9, 4.0, 3.2,
                         img_h=1.7, header_color=color)
    _rect(s, MARGIN, 5.25, SW - MARGIN * 2, 0.06, C_SILVER)
    _bullets(s, [
        "Im PI-Planning: historischen Durchsatz als Kapazitätsbasis nutzen (realistischer als Story Points).",
        "Wochenverlauf zeigt, ob Urlaube, Feiertage oder andere Faktoren den Durchsatz beeinflussen.",
        "Tipp: Filter auf bestimmte Issue-Typen setzen (z. B. nur 'Feature') für typ-spezifische Velocity.",
    ], MARGIN, 5.4, SW - MARGIN * 2, 1.8, size=12, color=C_TEXT, marker="💡", space=4)


def _r7_flow_load(prs):
    s = _blank(prs)
    _header(s, "Flow Load / WIP — Was läuft gerade?", color=C_ORANGE)
    _txt(s, "Messung: Alle Aufgaben, die aktiv begonnen wurden (First Date gesetzt), aber noch nicht abgeschlossen sind.",
         MARGIN, 1.42, SW - MARGIN * 2, 0.35, 13, C_TEXT)
    _card_with_image(s, "📦  Boxplot je Stage",
          IMG_METRICS / "Flow_load_AgingWork.png",
          "Y-Achse = Alter in Tagen (seit First Date).\n"
          "Referenzlinien aus historischen Abschlüssen: Median, P85, Target CT.\n"
          "Aufgaben über der P85-Linie: möglicherweise blockiert.",
          MARGIN, 1.88, 7.6, 3.2, img_h=1.7, header_color=C_ORANGE)
    _txt(s, "Was bedeutet es, wenn …", 8.2, 1.88, 4.8, 0.35, 13, C_DARK, bold=True)
    situations = [
        ("… viele Aufgaben in einer Stage alt werden:",
         "Engpass in dieser Phase — Kapazität prüfen."),
        ("… alle Stages voll sind:",
         "Zu viel WIP — neue Arbeit zurückhalten, erst fertigstellen."),
        ("… einzelne Aufgaben extrem alt sind:",
         "Blockade — konkrete Aufgabe prüfen."),
    ]
    for i, (q, a) in enumerate(situations):
        y = 2.3 + i * 0.95
        _rect(s, 8.2, y, 4.8, 0.82, C_LIGHT, line=C_SILVER)
        _txt(s, q, 8.35, y + 0.04, 4.5, 0.35, 11, C_DARK, bold=True)
        _txt(s, a, 8.35, y + 0.4,  4.5, 0.35, 11, C_GRAY)
    _rect(s, MARGIN, 5.25, SW - MARGIN * 2, 0.06, C_SILVER)
    _bullets(s, [
        "WIP-Limit: Je weniger Aufgaben gleichzeitig offen, desto schneller werden sie fertig (Little's Law).",
        "To-Do-Aufgaben (kein First Date) und abgeschlossene werden ausgeblendet — nur aktive Arbeit sichtbar.",
    ], MARGIN, 5.4, SW - MARGIN * 2, 1.8, size=12, color=C_TEXT, marker="💡", space=4)


def _r8_cfd(prs):
    s = _blank(prs)
    _header(s, "Cumulative Flow Diagram — Wie fließt die Arbeit?", color=C_GREEN)
    _txt(s, "Zeigt, wie viele Issues bis zu einem bestimmten Tag insgesamt in jede Stage eingetreten sind (kumuliert).",
         MARGIN, 1.42, SW - MARGIN * 2, 0.35, 13, C_TEXT)
    _card_with_image(s, "📉  Was man am CFD abliest",
          IMG_METRICS / "cummulative flow diagram.png",
          "Breite der Bänder = durchschnittliche Durchlaufzeit.\n"
          "→ Breiter: Prozess verlangsamt sich. Schmaler: Team wird schneller.\n"
          "• Oben (Inflow): Issues, die in die aktive Arbeit eintreten.\n"
          "• Unten (Outflow): Issues, die abgeschlossen werden.\n"
          "Inflow > Outflow über Zeit: Warteschlange wächst.",
          MARGIN, 1.88, 7.6, 4.6, img_h=2.5, header_color=C_GREEN)
    _txt(s, "Typische Muster:", 8.2, 1.88, 4.8, 0.35, 13, C_DARK, bold=True)
    patterns = [
        ("Gleichmäßige Bänder",  "Stabiler, vorhersagbarer Prozess.", C_GREEN),
        ("Breiter werdende Bänder", "Engpass aufgebaut — Outflow < Inflow.", C_ORANGE),
        ("Plateau (keine Bewegung)", "Work stoppage — Team blockiert.", C_RED),
        ("Schmaler werdende Bänder", "Aufholjagd — Rückstand wird abgebaut.", C_BLUE),
    ]
    for i, (pattern, desc, color) in enumerate(patterns):
        y = 2.3 + i * 1.0
        _rect(s, 8.2, y, 0.2, 0.75, color)
        _txt(s, pattern, 8.5, y + 0.03, 4.4, 0.35, 12, C_DARK, bold=True)
        _txt(s, desc,    8.5, y + 0.38, 4.4, 0.35, 11, C_GRAY)


def _r9_distribution(prs):
    s = _blank(prs)
    _header(s, "Flow Distribution — Was arbeitet das Team?", color=C_RED)
    _txt(s, "Drei Diagramme zeigen die Zusammensetzung der Arbeit.",
         MARGIN, 1.42, SW - MARGIN * 2, 0.35, 13, C_TEXT)
    charts = [
        ("🍩  By Issue Type",
         "Kreisdiagramm: Anteil von Feature,\nBug, Enabler usw.\n\n"
         "Frage: Wie viel Prozent unserer\nKapazität geht in Features,\nwie viel in Bugfixes?", C_RED),
        ("🍩  Stage Prominence",
         "Kreisdiagramm: In welcher Stage\nverbringen Aufgaben die meiste Zeit?\n\n"
         "Frage: Wo ist unser eigentlicher\nEngpass — auch wenn er nicht\noffensichtlich ist?", C_BLUE),
        ("📊  Avg Cycle Time by Type",
         "Balkendiagramm: Durchschnittliche\nDurchlaufzeit pro Aufgabentyp.\n\n"
         "Frage: Dauern Features wirklich\nlänger als Bugs? Wie groß ist\nder Unterschied?", C_ORANGE),
    ]
    for i, (title, body, color) in enumerate(charts):
        x = MARGIN + i * 4.25
        _card(s, title, body, x, 1.88, 4.0, 3.5, color)
    _rect(s, MARGIN, 5.55, SW - MARGIN * 2, 0.06, C_SILVER)
    _bullets(s, [
        "Stage Prominence zeigt, wo Aufgaben wirklich stecken — unabhängig von Tickets-Anzahl je Stage.",
        "Filter auf abgeschlossene Issues: Stage Prominence zählt nur Phasen vor dem Abschluss-Status.",
    ], MARGIN, 5.7, SW - MARGIN * 2, 1.6, size=12, color=C_TEXT, marker="💡", space=4)


def _r10_filter(prs):
    s = _blank(prs)
    _header(s, "Filter und Konfiguration", color=C_TEAL)
    _txt(s, "Filter:", MARGIN, 1.45, 3.0, 0.3, 14, C_DARK, bold=True)
    filter_rows = [
        ("Von / Bis", "Datumsfilter auf Closed Date — Standard: letzte 365 Tage"),
        ("Projekte", "Kommagetrennte Projekt-Keys, z. B.  ARTA, ARTB  — leer = alle"),
        ("Issuetypen", "z. B.  Feature, Bug  — leer = alle"),
        ("Status-Ausschluss", "Issues mit bestimmten Status komplett ignorieren, z. B.  Canceled"),
        ("Zero-Day-Issues", "Issues mit CT < Schwellwert (Standard: 5 min) ausschließen"),
    ]
    for i, (label, desc) in enumerate(filter_rows):
        y = 1.82 + i * 0.72
        _rect(s, MARGIN, y, SW - MARGIN * 2, 0.62, C_LIGHT, line=C_SILVER)
        _txt(s, label, MARGIN + 0.15, y + 0.08, 2.8, 0.35, 13, C_DARK, bold=True)
        _txt(s, desc,  MARGIN + 3.1,  y + 0.13, 9.8, 0.35, 12, C_GRAY)
    _rect(s, MARGIN, 5.5, SW - MARGIN * 2, 0.06, C_SILVER)
    _txt(s,
         "Tipp: einmal sauber konfigurierte Filter + Metrik-Auswahl als "
         "Project-Template speichern und in jedem Lauf wiederverwenden — "
         "siehe nächste Folie.",
         MARGIN, 5.7, SW - MARGIN * 2, 0.9, 13, C_BLUE, italic=True)


def _r_template(prs):
    """Slide: Shared project template (v0.11.0)."""
    s = _blank(prs)
    _header(s, "Project-Template — einmal speichern, in jedem Modul nutzen",
            color=C_TEAL)
    _txt(s, "Seit v0.11.0: eine einzige JSON-Datei mit einer Sektion pro Modul.",
         MARGIN, 1.42, SW - MARGIN * 2, 0.35, 13, C_TEXT)
    # JSON snippet (left)
    _rect(s, MARGIN, 1.88, 7.0, 4.4, C_LIGHT, line=C_SILVER)
    _txt(s, "Aufbau (schema v5):",
         MARGIN + 0.15, 1.95, 6.7, 0.35, 12, C_DARK, bold=True)
    json_lines = [
        "{",
        '  "schema": 5,',
        '  "transform_data":     { "json": "...", "workflow": "..." },',
        '  "build_reports":      { "pi": {...}, "filters": {...},',
        '                          "metrics": [...] },',
        '  "testdata_generator": { "out_dir": "...", "shape": "..." },',
        '  "helper":             { "inputs": [...] }',
        "}",
    ]
    _, tf = _txb(s, MARGIN + 0.15, 2.4, 6.7, 3.8)
    for i, line in enumerate(json_lines):
        _para(tf, line, 11, C_TEXT, first=(i == 0),
              space_before=2 if i > 0 else 0)
    # Bullet list (right)
    _txt(s, "Templates → Save / Load:",
         7.7, 1.88, 5.3, 0.35, 13, C_DARK, bold=True)
    _bullets(s, [
        ("Save in einem Modul: die anderen Modul-Sektionen "
         "bleiben unangetastet."),
        "Load in jedem Modul: füllt nur die eigene Sektion.",
        ("Hand-over (siehe Folie 5) nutzt denselben Template-"
         "Mechanismus als Transport."),
        "Alt-Format (schema v4) wird transparent als build_reports-Sektion gelesen.",
        "Ideal für regelmäßige Sprint-/PI-Berichte mit fixen Filtern.",
    ], 7.7, 2.3, 5.3, 4.0, size=12, color=C_TEXT, space=6)


def _r11_export(prs):
    s = _blank(prs)
    _header(s, "Export — Browser, PDF, Excel", color=C_BLUE)
    exports = [
        ("🌐  Browser-Anzeige",
         "Alle Diagramme als kombiniertes HTML im Browser.\n\n"
         "• Vollständig interaktiv: Zoom, Pan, Tooltip beim Hovern\n"
         "• Legende ein-/ausblenden\n"
         "• Kein PDF-Export nötig für interne Ansicht\n\n"
         "Aktivieren: Button 'Berichte erstellen' in der GUI\n"
         "CLI: --browser", C_BLUE),
        ("📄  PDF-Export",
         "Alle ausgewählten Metriken in einem\nmehrseitigen PDF.\n\n"
         "• Automatisch: Report-Excel-Datei (.xlsx)\n"
         "  mit allen gefilterten Issues\n"
         "• Automatisch: Zero-Day-Issues.xlsx\n"
         "  (wenn vorhanden)\n\n"
         "CLI: --pdf report.pdf", C_TEAL),
        ("📊  Report-Excel",
         "Wird bei jedem PDF-Export automatisch erstellt.\n\n"
         "Enthält alle gefilterten Issues mit:\n"
         "• Status Group (To Do / In Progress / Done)\n"
         "• Cycle Time Method A (Kalendertage)\n"
         "• Cycle Time Method B (Stage-Minuten)\n\n"
         "Ideal für eigene Auswertungen in Excel.", C_GREEN),
    ]
    for i, (title, body, color) in enumerate(exports):
        x = MARGIN + i * 4.25
        _card(s, title, body, x, 1.45, 4.05, 5.6, color)


def _r12_tips(prs):
    s = _blank(prs)
    _header(s, "Tipps für den Alltag", color=C_TEAL)
    tips = [
        ("Regelmäßiger Rhythmus",
         "Einmal pro Sprint oder PI ein festes Template laden und Report generieren. "
         "Konsistenz macht Trends erst sichtbar.", C_BLUE),
        ("Zero-Day-Issues bereinigen",
         "Beim ersten Lauf Zero-Day-Ausschluss aktivieren und prüfen, "
         "welche Issues herausgefiltert werden. Diese Aufgaben sind meist Testtickets "
         "oder Fehler im Jira-Workflow.", C_ORANGE),
        ("P85 als Service Level",
         "Den P85-Wert (85. Perzentil der Flow Time) als Service Level Expectation (SLE) "
         "definieren: '85 % unserer Features sind in X Tagen fertig.'", C_GREEN),
        ("CFD im Standup",
         "Das CFD eignet sich gut als Standup-Artefakt: Sehen wir Staubildung? "
         "Ist der Abfluss stabil? Veränderungen werden oft im CFD sichtbar, "
         "bevor sie in Tickets auffallen.", C_TEAL),
        ("Template-Pflege",
         "Einmal pro Sprint die Konfiguration als Project-Template speichern — "
         "Filter, PI-Konfig und Metrik-Auswahl bleiben dann konstant. "
         "Beim nächsten Lauf nur 'Templates → Laden' und 'Berichte erstellen'.", C_RED),
    ]
    for i, (title, body, color) in enumerate(tips):
        row, col = divmod(i, 2)
        if i == 4:  # last tip: full width
            x, w = MARGIN, SW - MARGIN * 2
        else:
            x = MARGIN + col * 6.4
            w = 6.2
        y = 1.45 + row * 2.0
        if i == 4:
            y = 1.45 + 2 * 2.0
        _card(s, "💡  " + title, body, x, y, w, 1.75, color)


def build_rte(path: Path) -> None:
    """Build the RTE/Scrum Master presentation and save to path."""
    prs = _new_prs()
    _r1_title(prs)
    _r_launcher(prs)
    _r2_workflow(prs)
    _r3_transform(prs)
    _r_handover(prs)
    _r4_build_intro(prs)
    _r5_flow_time(prs)
    _r6_flow_velocity(prs)
    _r7_flow_load(prs)
    _r8_cfd(prs)
    _r9_distribution(prs)
    _r10_filter(prs)
    _r_template(prs)
    _r11_export(prs)
    _r12_tips(prs)
    prs.save(path)
    print(f"RTE/SM PPTX gespeichert ({len(prs.slides)} Folien): {path.name}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    build_management(OUT_MGMT)
    build_rte(OUT_RTE)
    print("Fertig.")
