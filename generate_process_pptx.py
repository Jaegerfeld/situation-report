# -*- coding: utf-8 -*-
"""
Generate a business-process capability map for situation-report.
Target audience: non-technical stakeholders.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = Path(__file__).parent / "situation_report_prozess.pptx"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
C_BG     = RGBColor(0xf8, 0xf9, 0xfa)
C_DARK   = RGBColor(0x2c, 0x3e, 0x50)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_HINT   = RGBColor(0x9e, 0xa8, 0xb3)
C_BORDER = RGBColor(0xd5, 0xdb, 0xe1)

# One accent color per process step
STEP_COLORS = [
    RGBColor(0x29, 0x80, 0xb9),   # 1 – Daten bereitstellen  (blue)
    RGBColor(0x16, 0xa0, 0x85),   # 2 – Analyse einrichten   (teal)
    RGBColor(0x8e, 0x44, 0xad),   # 3 – Metriken berechnen   (purple)
    RGBColor(0x27, 0xae, 0x60),   # 4 – Ergebnisse teilen    (green)
    RGBColor(0xe6, 0x7e, 0x22),   # 5 – Konfiguration        (orange)
]

# ---------------------------------------------------------------------------
# Process definition
# ---------------------------------------------------------------------------

STEPS = [
    {
        "num": "1",
        "title": "Daten\nbereitstellen",
        "subtitle": "Rohdaten aus Jira aufbereiten",
        "icon": "📤",
        "tool": "transform_data",
        "short": [
            "Jira-Export importieren",
            "Workflow definieren",
            "Zeiten berechnen",
        ],
        "capabilities": [
            ("Jira-Daten importieren",
             "JSON-Export aus Jira einlesen -- kein manuelles Umformatieren noetig."),
            ("Workflow definieren",
             "Eine einfache Textdatei legt fest, welche Schritte zum Prozess gehoeren "
             "und wo Arbeit beginnt bzw. endet."),
            ("Verweildauer berechnen",
             "Das System ermittelt automatisch, wie lange jedes Ticket in jedem "
             "Prozessschritt verbracht hat."),
            ("Zeitpunkte bestimmen",
             "Beginn der Bearbeitung (First Date) und Abschluss (Closed Date) werden "
             "aus den Statuswechseln abgeleitet -- auch bei Wiedereroffnung korrekt."),
            ("Luecken behandeln",
             "Wird ein Prozessschritt uebersprungen, wird die Zeit automatisch "
             "dem naechsten bekannten Schritt zugerechnet."),
            ("Qualitaetsprufung",
             "Fehlende Workflow-Marker und unbekannte Status werden gemeldet, "
             "bevor Fehler in die Auswertung einfliessen."),
            ("Ausgabe als Excel",
             "Ergebnis: zwei Excel-Dateien (IssueTimes + CFD), die direkt von "
             "build_reports weiterverarbeitet werden."),
        ],
    },
    {
        "num": "2",
        "title": "Analyse\neinrichten",
        "subtitle": "Zeitraum, Projekte und Filter festlegen",
        "icon": "⚙️",
        "tool": "build_reports",
        "short": [
            "Zeitraum wahlen",
            "Projekte filtern",
            "Ausnahmen definieren",
        ],
        "capabilities": [
            ("Zeitraum festlegen",
             "Per Kalender Von/Bis wahlen oder auf Knopfdruck die letzten 365 Tage "
             "laden -- der Bericht passt sich sofort an."),
            ("Projekte auswahlen",
             "Bei mehreren Teams oder Projekten in einer Datei kann gezielt "
             "nur ein Projekt oder eine Kombination ausgewertet werden."),
            ("Issue-Typen filtern",
             "Bugs, Features, Stories getrennt oder gemeinsam betrachten."),
            ("Ausnahmen ausschliessen",
             "Tickets mit bestimmten Abschlussgrundon (z. B. 'Duplicate') oder "
             "Status (z. B. 'Canceled') werden vollstandig herausgerechnet."),
            ("Schnellausschluss Zero-Day",
             "Tickets, die in Sekunden durch den Workflow geklickt wurden, "
             "verzerren Statistiken -- sie werden automatisch erkannt und "
             "separat dokumentiert."),
            ("PI-Intervalle konfigurieren",
             "Eigene Program-Increment-Zeitraume definieren, damit Velocity-Berichte "
             "die echten PI-Grenzen des Teams zeigen."),
        ],
    },
    {
        "num": "3",
        "title": "Metriken\nberechnen",
        "subtitle": "Fuenf Flow-Metriken auf einen Blick",
        "icon": "📊",
        "tool": "build_reports",
        "short": [
            "Durchlaufzeit messen",
            "Durchsatz pruefen",
            "Engpaesse finden",
        ],
        "capabilities": [
            ("Wie lange dauert ein Ticket?\n(Flow Time / Cycle Time)",
             "Zeigt Durchlaufzeit-Verteilung, Trend ueber Zeit und statistische "
             "Kennzahlen -- auf einen Blick sichtbar, ob der Prozess stabiler oder "
             "langsamer wird."),
            ("Wie viel liefert das Team?\n(Flow Velocity / Throughput)",
             "Taegliche, wochentliche und PI-Lieferraten -- ideal fuer "
             "Kapazitatsplanung und Sprint-Reviews."),
            ("Wie viel Arbeit liegt parallel?\n(Flow Load / WIP)",
             "Zeigt, wie alt offene Tickets je Prozessschritt sind. "
             "Je mehr parallel, desto langsamer der Durchfluss."),
            ("Wie fliessen Tickets durchs System?\n(Cumulative Flow Diagram)",
             "Sichtbar, ob mehr eingeht als abgeschlossen wird -- und in welchem "
             "Schritt sich Tickets aufstauen."),
            ("Womit verbringt das Team Zeit?\n(Flow Distribution)",
             "Anteil Bugs vs. Features, welcher Schritt dominiert, "
             "welche Ticket-Typen am langsamsten sind."),
        ],
    },
    {
        "num": "4",
        "title": "Ergebnisse\nteilen",
        "subtitle": "Berichte erstellen und kommunizieren",
        "icon": "📋",
        "tool": "build_reports",
        "short": [
            "PDF exportieren",
            "Im Browser zeigen",
            "Excel bereitstellen",
        ],
        "capabilities": [
            ("Interaktive Browser-Ansicht",
             "Alle Diagramme im Browser: hineinzoomen, einzelne Werte per Hover "
             "ablesen, Kategorien ein-/ausblenden -- ideal fuer Live-Prasentationen."),
            ("PDF-Bericht",
             "Mehrseitiger PDF-Bericht mit allen Diagrammen -- versendbar, "
             "druckbar, archivierbar."),
            ("Report-Excel",
             "Jeder PDF-Export erzeugt automatisch eine Excel-Datei mit allen "
             "gefilterten Tickets, Statusgruppen und Durchlaufzeiten."),
            ("Zero-Day-Dokumentation",
             "Tickets, die als Ausreisser herausgefiltert wurden, erscheinen "
             "in einer separaten Excel-Datei -- transparent und nachvollziehbar."),
        ],
    },
    {
        "num": "5",
        "title": "Konfiguration\nverwalten",
        "subtitle": "Einstellungen sichern und Hilfe finden",
        "icon": "🔧",
        "tool": "beide Module",
        "short": [
            "Einstellungen speichern",
            "Sprache wahlen",
            "Handbuch offnen",
        ],
        "capabilities": [
            ("Einstellungen als Template sichern",
             "Alle Filter, Dateipfade und Optionen in einer JSON-Datei speichern "
             "und beim nachsten Mal mit einem Klick wieder laden."),
            ("Ausschlusse als Standard hinterlegen",
             "Haufig genutzte Ausschluss-Regeln dauerhaft speichern -- "
             "kein erneutes Eintippen bei jedem Start."),
            ("Sprache wahlen (DE / EN / RO / PT / FR)",
             "Launcher in 5 Sprachen nutzbar; build_reports und transform_data "
             "auf Deutsch oder Englisch. Die gewunschte Sprache wird beim "
             "nachsten Start beibehalten."),
            ("Terminologie anpassen",
             "SAFe-Begriffe (Flow Time, Flow Velocity ...) oder globale Begriffe "
             "(Cycle Time, Throughput ...) -- je nach Teamkonvention."),
            ("Benutzerhandbuch",
             "PDF-Handbuch: Launcher in 5 Sprachen (DE, EN, RO, PT, FR), "
             "build_reports und transform_data je DE + EN -- direkt aus dem "
             "Hilfemenu abrufbar."),
        ],
    },
]

# ---------------------------------------------------------------------------
# pptx helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill, line=None, line_w_pt=0.75, radius=False):
    from pptx.util import Inches
    shp = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w_pt)
    else:
        shp.line.fill.background()
    if radius:
        # rounded corners via XML
        sp = shp._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 30000')
    return shp


def txtbox(slide, text, l, t, w, h,
           size=10, bold=False, color=C_DARK,
           align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx


def add_paragraph(tf, text, size=10, bold=False, color=C_DARK,
                  align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


# ---------------------------------------------------------------------------
# Slide 1 – Title
# ---------------------------------------------------------------------------

def slide_title(prs):
    sl = blank_slide(prs)

    # Background
    rect(sl, 0, 0, 13.33, 7.5, C_DARK)

    # Accent band
    rect(sl, 0, 2.55, 13.33, 2.45, RGBColor(0x1a, 0x52, 0x76))

    # Color stripes for 5 steps at bottom
    stripe_w = 13.33 / 5
    for i, c in enumerate(STEP_COLORS):
        rect(sl, i * stripe_w, 7.0, stripe_w, 0.5, c)

    # Main title
    txtbox(sl, "situation-report",
           0.5, 2.7, 12.3, 1.1,
           size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    txtbox(sl, "Vom Jira-Export zum Entscheidungsdiagramm",
           0.5, 3.75, 12.3, 0.6,
           size=20, color=RGBColor(0xae, 0xcc, 0xe8), align=PP_ALIGN.CENTER)

    # Step labels in stripes
    for i, step in enumerate(STEPS):
        x = i * stripe_w + 0.08
        txtbox(sl, f"{step['num']}  {step['title'].replace(chr(10), ' ')}",
               x, 7.02, stripe_w - 0.1, 0.45,
               size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Footer
    txtbox(sl, "github.com/Jaegerfeld/situation-report  |  Stand April 2026",
           0.5, 6.55, 12.3, 0.35,
           size=9, color=C_HINT, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 – Process overview (all 5 steps + key capabilities)
# ---------------------------------------------------------------------------

def slide_overview(prs):
    sl = blank_slide(prs)

    # Background
    rect(sl, 0, 0, 13.33, 7.5, C_BG)

    # Header bar
    rect(sl, 0, 0, 13.33, 0.85, C_DARK)
    txtbox(sl, "Geschaftsprozess: Von Rohdaten zur Entscheidungsgrundlage",
           0.3, 0.1, 12.7, 0.65,
           size=18, bold=True, color=C_WHITE)

    # Step boxes
    box_w   = 2.22
    box_h   = 1.4
    gap     = 0.14
    total_w = 5 * box_w + 4 * gap
    x0      = (13.33 - total_w) / 2
    box_y   = 1.1

    for i, step in enumerate(STEPS):
        c   = STEP_COLORS[i]
        lc  = RGBColor(int(c[0]*0.75), int(c[1]*0.75), int(c[2]*0.75))
        x   = x0 + i * (box_w + gap)

        # Colored box
        rect(sl, x, box_y, box_w, box_h, c, radius=True)

        # Number badge
        rect(sl, x + 0.1, box_y + 0.08, 0.38, 0.38,
             RGBColor(0xff, 0xff, 0xff), radius=True)
        txtbox(sl, step["num"],
               x + 0.1, box_y + 0.06, 0.38, 0.38,
               size=13, bold=True, color=c, align=PP_ALIGN.CENTER)

        # Icon + title
        txtbox(sl, step["icon"],
               x + 0.55, box_y + 0.05, 0.45, 0.45,
               size=18, align=PP_ALIGN.CENTER)
        txtbox(sl, step["title"],
               x + 0.1, box_y + 0.58, box_w - 0.2, 0.75,
               size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < 4:
            arr_x = x + box_w + 0.01
            txtbox(sl, "▶",
                   arr_x, box_y + 0.5, gap + 0.1, 0.45,
                   size=11, color=C_HINT, align=PP_ALIGN.CENTER)

        # Capability bullets below box
        cap_y = box_y + box_h + 0.12
        rect(sl, x, cap_y, box_w, 0.025, c)   # thin color line

        for j, bullet in enumerate(step["short"]):
            by = cap_y + 0.05 + j * 0.42
            rect(sl, x + 0.08, by + 0.1, 0.12, 0.12, c, radius=True)
            txtbox(sl, bullet,
                   x + 0.26, by + 0.01, box_w - 0.32, 0.38,
                   size=8.5, color=C_DARK)

    # Module labels
    module_y = 6.5
    rect(sl, 0.35, module_y, 4.5, 0.55, RGBColor(0xe8, 0xf4, 0xfb))
    txtbox(sl, "Modul: transform_data",
           0.45, module_y + 0.05, 4.3, 0.45,
           size=8.5, bold=True, color=STEP_COLORS[0], align=PP_ALIGN.CENTER)
    txtbox(sl, "Schritt 1",
           0.45, module_y + 0.24, 4.3, 0.25,
           size=7.5, color=C_HINT, align=PP_ALIGN.CENTER)

    rect(sl, 5.1, module_y, 7.9, 0.55, RGBColor(0xf0, 0xf0, 0xf8))
    txtbox(sl, "Modul: build_reports",
           5.2, module_y + 0.05, 7.7, 0.45,
           size=8.5, bold=True, color=STEP_COLORS[2], align=PP_ALIGN.CENTER)
    txtbox(sl, "Schritte 2 - 5",
           5.2, module_y + 0.24, 7.7, 0.25,
           size=7.5, color=C_HINT, align=PP_ALIGN.CENTER)

    # Footer
    rect(sl, 0, 7.3, 13.33, 0.2, C_DARK)
    txtbox(sl, "situation-report  |  Flow-Metriken fur agile Teams",
           0.3, 7.3, 12.7, 0.18,
           size=7.5, color=C_HINT, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 3-7 – One detail slide per step
# ---------------------------------------------------------------------------

def slide_step(prs, step_idx):
    step  = STEPS[step_idx]
    color = STEP_COLORS[step_idx]
    sl    = blank_slide(prs)
    caps  = step["capabilities"]

    # Background
    rect(sl, 0, 0, 13.33, 7.5, C_BG)

    # Header
    rect(sl, 0, 0, 13.33, 1.15, color)
    # Step number pill
    rect(sl, 0.25, 0.18, 0.55, 0.55, C_WHITE, radius=True)
    txtbox(sl, step["num"],
           0.25, 0.15, 0.55, 0.55,
           size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Icon
    txtbox(sl, step["icon"],
           0.9, 0.14, 0.55, 0.55,
           size=22, align=PP_ALIGN.LEFT)
    # Title + subtitle
    txtbox(sl, step["title"].replace("\n", " "),
           1.55, 0.1, 8.0, 0.55,
           size=22, bold=True, color=C_WHITE)
    txtbox(sl, step["subtitle"],
           1.55, 0.62, 8.0, 0.42,
           size=11, color=RGBColor(0xff, 0xff, 0xff), italic=True)
    # Module badge
    rect(sl, 10.2, 0.25, 2.85, 0.45,
         RGBColor(0xff, 0xff, 0xff))
    txtbox(sl, f"Modul: {step['tool']}",
           10.25, 0.27, 2.75, 0.38,
           size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Process ribbon (mini overview, highlight current)
    ribbon_y = 1.25
    ribbon_h = 0.48
    box_w_sm = 2.35
    gap_sm   = 0.04
    total_sm = 5 * box_w_sm + 4 * gap_sm
    rx0      = (13.33 - total_sm) / 2

    for i, s in enumerate(STEPS):
        is_cur = (i == step_idx)
        c      = STEP_COLORS[i] if is_cur else RGBColor(0xcc, 0xd6, 0xdd)
        x      = rx0 + i * (box_w_sm + gap_sm)
        rect(sl, x, ribbon_y, box_w_sm, ribbon_h, c, radius=True)
        fc = C_WHITE if is_cur else C_HINT
        txtbox(sl, f"{s['num']}  {s['title'].replace(chr(10), ' ')}",
               x + 0.05, ribbon_y + 0.06, box_w_sm - 0.1, ribbon_h - 0.1,
               size=8 if is_cur else 7.5,
               bold=is_cur, color=fc, align=PP_ALIGN.CENTER)

    # Capability cards grid
    n      = len(caps)
    cols   = 3 if n > 4 else 2
    rows   = -(-n // cols)   # ceiling division
    card_w = (13.33 - 0.6) / cols
    card_h = (7.5 - 1.95) / rows - 0.1
    card_h = min(card_h, 1.55)

    for idx, (cap_title, cap_desc) in enumerate(caps):
        col = idx % cols
        row = idx // cols
        cx  = 0.3 + col * card_w
        cy  = 1.95 + row * (card_h + 0.1)
        cw  = card_w - 0.15

        # Card background
        rect(sl, cx, cy, cw, card_h, C_WHITE, C_BORDER, line_w_pt=0.5)
        # Colored left accent bar
        rect(sl, cx, cy, 0.07, card_h, color)

        # Title
        txtbox(sl, cap_title,
               cx + 0.15, cy + 0.1, cw - 0.22, 0.38,
               size=9.5, bold=True, color=color)
        # Description
        txtbox(sl, cap_desc,
               cx + 0.15, cy + 0.43, cw - 0.22, card_h - 0.5,
               size=8.5, color=C_DARK)

    # Footer
    rect(sl, 0, 7.35, 13.33, 0.15, color)
    txtbox(sl, "situation-report  |  Flow-Metriken fur agile Teams",
           0.3, 7.35, 12.7, 0.13,
           size=7, color=C_WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    slide_title(prs)
    slide_overview(prs)
    for i in range(len(STEPS)):
        slide_step(prs, i)

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
