"""Generate a PowerPoint feature overview for situation-report."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
C_DARK   = RGBColor(0x2c, 0x3e, 0x50)   # dark blue – headings / title bg
C_ACCENT = RGBColor(0x29, 0x80, 0xb9)   # mid blue  – section headers
C_LIGHT  = RGBColor(0xec, 0xf0, 0xf1)   # light grey – row stripes / bg
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_HINT   = RGBColor(0x7f, 0x8c, 0x8d)   # grey      – captions
C_GREEN  = RGBColor(0x27, 0xae, 0x60)   # infra slide accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUTPUT = Path(__file__).parent / "situation_report_features.pptx"

# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------

SLIDES = [
    # ── Title slide ─────────────────────────────────────────────────────────
    {
        "type": "title",
        "title": "situation-report",
        "subtitle": "Feature-Übersicht — Stand April 2026",
    },
    # ── transform_data ──────────────────────────────────────────────────────
    {
        "type": "table",
        "heading": "transform_data",
        "rows": [
            ("Kern-Transformation",
             "Liest Jira-JSON-Export, berechnet Verweildauer je Workflow-Stage, "
             "schreibt IssueTimes.xlsx und CFD.xlsx"),
            ("Workflow-Datei",
             "Textdatei mit <First>- und <Closed>-Markern definiert Eintritts- "
             "und Abschluss-Stage"),
            ("Datumsregeln",
             "First Date / Closed Date aus Stage-Übergängen; Fallback bei übersprungenen "
             "Stages; Closed Date wird gelöscht wenn Issue wiedereröffnet wurde"),
            ("Carry-Forward",
             "Verweildauer wird auf letzte bekannte Stage weitergezählt wenn eine "
             "Stage übersprungen wurde"),
            ("Validierung & Warnungen",
             "Warnt bei fehlenden Workflow-Markern, unmapped Statuses und "
             "Stage-Inkonsistenzen"),
            ("GUI",
             "tkinter-Oberfläche mit Dateidialogen für JSON, Workflow-Datei und "
             "Ausgabeordner; Präfix automatisch aus Dateiname vorbelegt"),
            ("Ladebalken",
             "Erscheint nach 3 Sekunden wenn die Transformation noch läuft"),
            ("Doppelklick-Starter",
             "transform_data_gui.pyw startet die GUI ohne Konsolenfenster"),
        ],
    },
    # ── build_reports: Metriken ──────────────────────────────────────────────
    {
        "type": "table",
        "heading": "build_reports — Metriken",
        "rows": [
            ("Flow Time / Cycle Time",
             "Boxplot + Scatterplot; Methode A & B; LOESS-Trendlinie; "
             "Perzentil-Referenzlinien; farbige Punkte; rote Ausreißer; "
             "Issue-Key als Hover-Tooltip"),
            ("Flow Velocity / Throughput",
             "Tagesfrequenz-Histogramm + Wochenverlauf-Linienchart + "
             "PI-Balkendiagramm; konfigurierbares Target CT"),
            ("Flow Load / WIP",
             "Aging-WIP-Boxplot pro Stage; Issue-Anzahl je Stage als "
             "Annotation; Referenzlinien (Median, P85, Target CT); Legende"),
            ("CFD",
             "Gestapeltes Flächendiagramm; Trendlinien an visuellen "
             "Oberkanten der <First>- und <Closed>-Stages; "
             "Monats-/KW-Achse ohne Überlappung; In/Out-Ratio im Titel"),
            ("Flow Distribution",
             "Drei Teildiagramme: Issue-Typ-Donut, Stage-Prominence-Donut, "
             "Ø Cycle Time je Typ als Balken"),
        ],
    },
    # ── build_reports: Filter ────────────────────────────────────────────────
    {
        "type": "table",
        "heading": "build_reports — Filter & Ausschlüsse",
        "rows": [
            ("Datumsfilter",
             'Von/Bis mit Kalender-Picker; "Letzte 365 Tage"-Button'),
            ("Projektfilter",
             "Auswahl aus allen in IssueTimes vorhandenen Projekten"),
            ("Issuetype-Filter",
             "Auswahl aus allen vorhandenen Issue-Typen"),
            ("Ausschluss Status / Resolution",
             "Issues komplett aus allen Metriken entfernen"),
            ("Zero-Day-Filter",
             "Konfigurierbarer Schwellwert (Standard 5 min); ausgeschlossene "
             "Issues in separater Excel dokumentiert"),
            ("Datumsfilter-Bugfix",
             "Offene Issues (kein Closed Date) wurden fälschlich vom "
             "Datumsfilter ausgeschlossen — behoben"),
        ],
    },
    # ── build_reports: Export ────────────────────────────────────────────────
    {
        "type": "table",
        "heading": "build_reports — Export & Reporting",
        "rows": [
            ("Browser-Anzeige",
             "Vollständig interaktive Plotly-Diagramme im Browser"),
            ("PDF-Export",
             "Mehrseitige PDF mit allen gewählten Metriken"),
            ("Report-Excel",
             "Wird automatisch beim PDF-Export erzeugt; enthält gefilterte "
             "Issues mit Status Group, CT Methode A und B"),
            ("PI-Konfiguration",
             "Eigene PI-Intervalle per JSON-Datei; ohne Datei werden "
             "Kalenderquartale verwendet"),
        ],
    },
    # ── build_reports: GUI ───────────────────────────────────────────────────
    {
        "type": "table",
        "heading": "build_reports — GUI & UX",
        "rows": [
            ("Terminologie",
             "Umschaltung SAFe ↔ Global (Flow Time ↔ Cycle Time etc.)"),
            ("Templates",
             "Alle Einstellungen als JSON speichern/laden; Ausschlüsse als "
             "dauerhafter Standard speicherbar"),
            ("Tooltips",
             "Hover-Tooltips auf allen GUI-Elementen"),
            ("Ladebalken",
             "Erscheint nach 3 Sekunden bei laufenden Berechnungen"),
            ("Stage-Konsistenzprüfung",
             "Warnung wenn IssueTimes und CFD unterschiedliche Stages haben"),
            ("Kollisionsfreie Annotations",
             "Referenzlinien-Labels weichen automatisch aus wenn sie "
             "sich überlappen"),
            ("Doppelklick-Starter",
             "build_reports_gui.pyw startet die GUI ohne Konsolenfenster"),
        ],
    },
    # ── Infrastruktur ────────────────────────────────────────────────────────
    {
        "type": "table",
        "heading": "Infrastruktur & Dokumentation",
        "accent": C_GREEN,
        "rows": [
            ("Versionierung",
             "SemVer ab v0.2.0; zentrale version.py wird von beiden GUIs "
             "und den Manuals gelesen"),
            ("Sprachauswahl",
             "DE/EN in beiden GUIs; Flagge 🇩🇪/🇬🇧 ganz rechts im Menüband; "
             "letzte Wahl wird in ~/.situation_report/prefs.json gespeichert"),
            ("PDF-Manuals",
             "Benutzerhandbuch (DE) + User Manual (EN) für build_reports "
             "und transform_data; generiert mit ReportLab inkl. echter "
             "Beispieldiagramme"),
            ("Manual-Link in GUI",
             "Hilfe-Menü öffnet das sprachpassende PDF direkt auf GitHub Pages"),
            ("GitHub Pages Doku",
             "MkDocs-Site mit DE/EN-Doku für alle Module; "
             "Architektur nach C4-Modell"),
            ("CI/CD",
             "GitHub Actions für automatisches Docs-Deployment"),
            ("Tests",
             "Unit- und Acceptance-Tests für beide Module; aktuell 448 Tests"),
        ],
    },
]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=11, bold=False, color=C_DARK,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title  # may be None on blank layout — ignore

    # Full-bleed dark background
    add_rect(slide, 0, 0, 13.33, 7.5, C_DARK)
    # Accent band
    add_rect(slide, 0, 2.6, 13.33, 2.4, C_ACCENT)

    # Title
    add_text_box(slide, data["title"],
                 0.5, 2.8, 12.3, 1.2,
                 font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, data["subtitle"],
                 0.5, 3.85, 12.3, 0.7,
                 font_size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    # Footer
    add_text_box(slide, "github.com/Jaegerfeld/situation-report",
                 0.5, 6.8, 12.3, 0.4,
                 font_size=9, color=C_HINT, align=PP_ALIGN.CENTER)


def build_table_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    accent = data.get("accent", C_ACCENT)
    rows = data["rows"]

    # Header bar
    add_rect(slide, 0, 0, 13.33, 0.9, C_DARK)
    add_text_box(slide, data["heading"],
                 0.3, 0.1, 12.7, 0.7,
                 font_size=22, bold=True, color=C_WHITE)

    # Column headers
    add_rect(slide, 0.3, 1.0, 3.2, 0.35, C_ACCENT)
    add_rect(slide, 3.6, 1.0, 9.4, 0.35, C_ACCENT)
    add_text_box(slide, "Feature", 0.35, 1.02, 3.1, 0.3,
                 font_size=9, bold=True, color=C_WHITE)
    add_text_box(slide, "Beschreibung", 3.65, 1.02, 9.3, 0.3,
                 font_size=9, bold=True, color=C_WHITE)

    # Rows
    row_h = (7.5 - 1.45) / max(len(rows), 1)
    row_h = min(row_h, 0.72)
    y = 1.45

    for i, (feature, desc) in enumerate(rows):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        add_rect(slide, 0.3, y, 3.2, row_h, bg)
        add_rect(slide, 3.6, y, 9.4, row_h, bg)

        add_text_box(slide, feature, 0.38, y + 0.04, 3.05, row_h - 0.06,
                     font_size=9, bold=True, color=C_DARK)
        add_text_box(slide, desc, 3.68, y + 0.04, 9.2, row_h - 0.06,
                     font_size=9, color=C_DARK)
        y += row_h

    # Footer line
    add_rect(slide, 0, 7.35, 13.33, 0.15, accent)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for data in SLIDES:
        if data["type"] == "title":
            build_title_slide(prs, data)
        elif data["type"] == "table":
            build_table_slide(prs, data)

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
